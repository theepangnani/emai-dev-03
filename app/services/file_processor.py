"""
File processor service for extracting text from various document formats.
Supports: PDF, Word, Excel, PowerPoint, Images (OCR), Text, and ZIP archives.
"""

import base64
import io
import os
import zipfile
import tempfile
from pathlib import Path
from typing import BinaryIO

from app.core.logging_config import get_logger
from app.core.config import settings

# Document processing
import PyPDF2
from docx import Document as WordDocument
from pptx import Presentation
from openpyxl import load_workbook
from PIL import Image

# OCR (optional - requires Tesseract installed)
try:
    import pytesseract
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

# Vision-based OCR via Anthropic Claude (preferred for math/formulas)
try:
    import anthropic
    VISION_OCR_AVAILABLE = True
except ImportError:
    VISION_OCR_AVAILABLE = False

# PDF-to-image conversion (optional - requires poppler-utils)
try:
    from pdf2image import convert_from_bytes
    PDF2IMAGE_AVAILABLE = True
except ImportError:
    PDF2IMAGE_AVAILABLE = False

# Constants — size limit driven by settings so it's configurable per environment
MAX_FILE_SIZE = settings.max_upload_size_mb * 1024 * 1024
SUPPORTED_EXTENSIONS = {
    '.pdf', '.doc', '.docx', '.txt', '.md', '.rtf',
    '.xlsx', '.xls', '.csv',
    '.pptx', '.ppt',
    '.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp',
    '.zip'
}

IMAGE_EXTENSIONS = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp'}

# Minimum extracted text length to proceed with AI generation (#2217)
MIN_EXTRACTED_TEXT_LENGTH = 50
MIN_EXTRACTION_CHARS = MIN_EXTRACTED_TEXT_LENGTH  # Alias for import convenience

logger = get_logger(__name__)

def check_extracted_text_sufficient(text: str, filename: str = "unknown") -> None:
    """Raise FileProcessingError if extracted text is too short for AI generation.

    Args:
        text: The extracted text content.
        filename: Original filename for logging context.

    Raises:
        FileProcessingError: When stripped text length is below MIN_EXTRACTED_TEXT_LENGTH.
    """
    stripped = text.strip() if text else ""
    if len(stripped) < MIN_EXTRACTED_TEXT_LENGTH:
        logger.warning(
            "Blocked AI generation — insufficient text extracted. "
            "file=%s, chars=%d, threshold=%d",
            filename,
            len(stripped),
            MIN_EXTRACTED_TEXT_LENGTH,
        )
        raise FileProcessingError(
            "We couldn't read this document. Try uploading a text-based PDF "
            "or use the text input option."
        )


# Magic bytes signatures for binary formats.
# Each entry maps an extension to a list of accepted prefixes (bytes).
# Text-based formats (.txt, .md, .csv, .rtf) are omitted — content can start arbitrarily.
_MAGIC_BYTES: dict[str, list[bytes]] = {
    '.pdf':  [b'%PDF-'],
    '.png':  [b'\x89PNG\r\n\x1a\n'],
    '.jpg':  [b'\xff\xd8\xff'],
    '.jpeg': [b'\xff\xd8\xff'],
    '.gif':  [b'GIF87a', b'GIF89a'],
    '.bmp':  [b'BM'],
    '.tiff': [b'II*\x00', b'MM\x00*'],
    '.webp': [b'RIFF'],
    # ZIP-based (docx, xlsx, pptx, zip all start with PK)
    '.zip':  [b'PK\x03\x04', b'PK\x05\x06', b'PK\x07\x08'],
    '.docx': [b'PK\x03\x04'],
    '.xlsx': [b'PK\x03\x04'],
    '.pptx': [b'PK\x03\x04'],
    # Legacy OLE2 binary Office formats (doc, xls, ppt)
    '.doc':  [b'\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1'],
    '.xls':  [b'\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1'],
    '.ppt':  [b'\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1'],
}


class FileProcessingError(Exception):
    """Custom exception for file processing errors."""
    pass


def _check_magic_bytes(file_content: bytes, ext: str) -> bool:
    """Return True if file_content starts with a known-good magic prefix for ext.

    Returns True unconditionally for extensions not in _MAGIC_BYTES (text formats).
    """
    expected = _MAGIC_BYTES.get(ext)
    if not expected:
        return True  # no magic bytes rule for this type
    return any(file_content.startswith(sig) for sig in expected)


def validate_file(file_content: bytes, filename: str) -> None:
    """Validate file size, extension allowlist, and magic bytes."""
    file_size_mb = len(file_content) / (1024 * 1024)
    logger.debug(f"Validating file: {filename}, size: {file_size_mb:.2f} MB")

    if len(file_content) > MAX_FILE_SIZE:
        logger.warning(f"File too large: {filename} ({file_size_mb:.2f} MB)")
        raise FileProcessingError(
            f"File size exceeds maximum allowed size of {MAX_FILE_SIZE // (1024*1024)} MB"
        )

    ext = Path(filename).suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        logger.warning(f"Unsupported file type: {ext} for file {filename}")
        raise FileProcessingError(
            f"Unsupported file type: {ext}. Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )

    if not _check_magic_bytes(file_content, ext):
        logger.warning(f"Magic bytes mismatch for {filename} (claimed ext: {ext})")
        raise FileProcessingError(
            f"File content does not match the expected format for {ext} files."
        )


def extract_text_from_pdf(file_content: bytes) -> str:
    """Extract text from PDF file. Falls back to OCR for scanned/image PDFs."""
    try:
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
        text_parts = []
        empty_pages = 0
        page_count = len(pdf_reader.pages)
        logger.debug(f"Processing PDF with {page_count} pages")
        for page in pdf_reader.pages:
            text = page.extract_text()
            if text and text.strip():
                text_parts.append(text)
            else:
                empty_pages += 1

        # If most pages had no text, try OCR fallback for scanned PDFs
        if empty_pages > len(text_parts) and OCR_AVAILABLE and PDF2IMAGE_AVAILABLE:
            logger.info(f"PDF has {empty_pages} empty pages, attempting OCR fallback")
            try:
                images = convert_from_bytes(file_content, dpi=200)
                ocr_parts = []
                for i, img in enumerate(images):
                    ocr_text = pytesseract.image_to_string(img)
                    if ocr_text.strip():
                        ocr_parts.append(ocr_text.strip())
                if ocr_parts:
                    logger.info(f"OCR extracted text from {len(ocr_parts)}/{len(images)} pages")
                    return "\n\n".join(ocr_parts)
            except Exception as ocr_err:
                logger.warning(f"PDF OCR fallback failed: {ocr_err}")

        logger.debug(f"Extracted text from {len(text_parts)} pages")
        return "\n\n".join(text_parts)
    except Exception as e:
        logger.error(f"PDF extraction failed: {str(e)}")
        raise FileProcessingError(f"Failed to extract text from PDF: {str(e)}")


def _extract_text_from_docx_textboxes(doc) -> list[str]:
    """Extract text from text boxes and shapes in a Word document via XML."""
    text_parts = []
    nsmap = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    for txbx in doc.element.body.iter(
        "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}txbxContent"
    ):
        for paragraph in txbx.findall(".//w:p", nsmap):
            runs = paragraph.findall(".//w:r/w:t", nsmap)
            para_text = "".join(r.text or "" for r in runs)
            if para_text.strip():
                text_parts.append(para_text.strip())
    return text_parts


def _extract_images_from_docx(file_content: bytes) -> list[bytes]:
    """Extract embedded images from a .docx file (which is a ZIP archive)."""
    images = []
    try:
        with zipfile.ZipFile(io.BytesIO(file_content), "r") as zf:
            for name in zf.namelist():
                if name.startswith("word/media/") and Path(name).suffix.lower() in IMAGE_EXTENSIONS:
                    images.append(zf.read(name))
    except Exception as e:
        logger.debug(f"Failed to extract images from docx: {e}")
    return images


def _get_image_media_type(img_bytes: bytes) -> str:
    """Detect image media type from file header bytes."""
    if img_bytes[:8] == b'\x89PNG\r\n\x1a\n':
        return "image/png"
    if img_bytes[:2] == b'\xff\xd8':
        return "image/jpeg"
    if img_bytes[:4] == b'GIF8':
        return "image/gif"
    if img_bytes[:2] in (b'BM',):
        return "image/bmp"
    # Default to png for unknown
    return "image/png"


def _ocr_images_with_vision(images: list[bytes], batch_size: int = 10) -> list[str]:
    """
    Use Claude Vision API to extract text from images (much better for math/formulas).

    Returns one description string per input image (empty string for skipped/tiny images).
    Batches images into groups to minimize API calls while staying within limits.
    Falls back gracefully if the API is unavailable or fails.
    """
    from app.core.config import settings

    if not VISION_OCR_AVAILABLE or not settings.anthropic_api_key:
        logger.debug("Vision OCR not available (no API key or anthropic not installed)")
        return [""] * len(images)

    # Build per-image results, default to empty string
    results: list[str] = [""] * len(images)

    # Filter out very small images (likely icons/decorators, < 1KB)
    meaningful_images = [(i, img) for i, img in enumerate(images) if len(img) > 1024]
    if not meaningful_images:
        return results

    logger.info(
        f"Vision OCR: processing {len(meaningful_images)}/{len(images)} images "
        f"(skipped {len(images) - len(meaningful_images)} tiny images)"
    )

    client = anthropic.Anthropic(api_key=settings.anthropic_api_key)
    # Use Haiku for OCR — fast, cheap, and accurate enough for text extraction
    ocr_model = "claude-haiku-4-5-20251001"

    # Process in batches
    for batch_start in range(0, len(meaningful_images), batch_size):
        batch = meaningful_images[batch_start:batch_start + batch_size]
        content_blocks = []

        if len(batch) == 1:
            # Single image — simpler prompt, no separator needed
            content_blocks.append({
                "type": "text",
                "text": (
                    "Extract ALL text from the following image. It is from an educational document "
                    "and may contain math formulas, equations, graphs, diagrams with labels, or handwritten notes.\n\n"
                    "IMPORTANT instructions:\n"
                    "- For math formulas, use clear notation: √ for square roots, ² for squared, "
                    "x₁ for subscripts, fractions as (a/b), etc.\n"
                    "- Preserve the structure and order of the content\n"
                    "- If the image is a graph or diagram, describe what it shows and include any visible labels/values\n"
                    "- If the image contains no meaningful text (e.g., a decorative element or icon), respond with [no text]\n"
                    "- Output ONLY the extracted text, no commentary"
                ),
            })
        else:
            # Multiple images — request separator between each
            content_blocks.append({
                "type": "text",
                "text": (
                    f"Extract ALL text from each of the following {len(batch)} images. They are from an educational document "
                    "and may contain math formulas, equations, graphs, diagrams with labels, or handwritten notes.\n\n"
                    "IMPORTANT instructions:\n"
                    "- For math formulas, use clear notation: √ for square roots, ² for squared, "
                    "x₁ for subscripts, fractions as (a/b), etc.\n"
                    "- Preserve the structure and order of the content\n"
                    "- If an image is a graph or diagram, describe what it shows and include any visible labels/values\n"
                    "- If an image contains no meaningful text (e.g., a decorative element or icon), respond with [no text]\n"
                    "- IMPORTANT: Separate the output for each image with the exact line: ---IMAGE_SEP---\n"
                    f"- You MUST output exactly {len(batch) - 1} separator lines (one between each image's text)\n"
                    "- Output ONLY the extracted text with separators, no commentary"
                ),
            })

        for idx, img_bytes in batch:
            media_type = _get_image_media_type(img_bytes)
            b64_data = base64.b64encode(img_bytes).decode("utf-8")
            content_blocks.append({
                "type": "image",
                "source": {
                    "type": "base64",
                    "media_type": media_type,
                    "data": b64_data,
                },
            })

        try:
            response = client.messages.create(
                model=ocr_model,
                max_tokens=4096,
                messages=[{"role": "user", "content": content_blocks}],
                temperature=0.0,
            )
            result_text = response.content[0].text.strip()
            logger.debug(
                f"Vision OCR batch {batch_start // batch_size + 1}: "
                f"extracted {len(result_text)} chars from {len(batch)} images"
            )

            if len(batch) == 1:
                # Single image result
                orig_idx = batch[0][0]
                if result_text and result_text != "[no text]":
                    results[orig_idx] = result_text
            else:
                # Split by separator for multi-image batches
                parts = result_text.split("---IMAGE_SEP---")
                for j, (orig_idx, _) in enumerate(batch):
                    if j < len(parts):
                        part = parts[j].strip()
                        if part and part != "[no text]":
                            results[orig_idx] = part
        except Exception as e:
            logger.warning(f"Vision OCR batch failed: {e}")
            # Don't abort — continue with remaining batches

    return results


def extract_text_from_docx(file_content: bytes) -> str:
    """Extract text from Word document (.docx).

    Extracts from paragraphs, tables, text boxes/shapes, and headers/footers.
    Falls back to OCR on embedded images if insufficient text is found.
    """
    try:
        doc = WordDocument(io.BytesIO(file_content))
        text_parts = []

        # 1. Extract from paragraphs
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)

        # 2. Extract from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    text_parts.append(" | ".join(row_text))

        # 3. Extract from text boxes / shapes (common in docs with images)
        try:
            textbox_parts = _extract_text_from_docx_textboxes(doc)
            text_parts.extend(textbox_parts)
            if textbox_parts:
                logger.debug(f"Extracted {len(textbox_parts)} text segments from text boxes")
        except Exception as e:
            logger.debug(f"Text box extraction skipped: {e}")

        # 4. Extract from headers and footers
        for section in doc.sections:
            for header in [section.header, section.first_page_header, section.even_page_header]:
                if header and header.is_linked_to_previous is False:
                    for p in header.paragraphs:
                        if p.text.strip():
                            text_parts.append(p.text)
            for footer in [section.footer, section.first_page_footer, section.even_page_footer]:
                if footer and footer.is_linked_to_previous is False:
                    for p in footer.paragraphs:
                        if p.text.strip():
                            text_parts.append(p.text)

        # 5. OCR on all embedded images
        # Always OCR embedded images — screenshots often contain the real content
        # (e.g., math problems, diagrams with text, scanned worksheets)
        images = _extract_images_from_docx(file_content)
        if images:
            total_text_len = sum(len(t) for t in text_parts)
            logger.info(
                f"Docx has {len(images)} embedded images ({total_text_len} chars text)"
            )

            # Prefer Vision OCR (Claude) — much better at math, formulas, diagrams
            vision_parts = _ocr_images_with_vision(images)
            non_empty_parts = [p for p in vision_parts if p]
            if non_empty_parts:
                logger.info(f"Vision OCR extracted {len(non_empty_parts)} text segments from images")
                text_parts.extend(non_empty_parts)
            elif OCR_AVAILABLE:
                # Fallback to Tesseract if Vision OCR unavailable/failed
                logger.info("Falling back to Tesseract OCR for embedded images")
                ocr_parts = []
                for i, img_bytes in enumerate(images):
                    try:
                        img = Image.open(io.BytesIO(img_bytes))
                        if img.mode in ("RGBA", "LA", "P"):
                            img = img.convert("RGB")
                        ocr_text = pytesseract.image_to_string(img)
                        if ocr_text.strip():
                            ocr_parts.append(ocr_text.strip())
                    except Exception as ocr_err:
                        logger.debug(f"OCR failed on embedded image {i}: {ocr_err}")
                if ocr_parts:
                    logger.info(f"Tesseract OCR extracted text from {len(ocr_parts)}/{len(images)} images")
                    text_parts.extend(ocr_parts)

        return "\n\n".join(text_parts)
    except Exception as e:
        if isinstance(e, FileProcessingError):
            raise
        raise FileProcessingError(f"Failed to extract text from Word document: {str(e)}")


def extract_text_from_pptx(file_content: bytes) -> str:
    """Extract text from PowerPoint presentation (.pptx)."""
    try:
        prs = Presentation(io.BytesIO(file_content))
        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"--- Slide {slide_num} ---"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            if len(slide_text) > 1:  # More than just the slide header
                text_parts.append("\n".join(slide_text))
        return "\n\n".join(text_parts)
    except Exception as e:
        raise FileProcessingError(f"Failed to extract text from PowerPoint: {str(e)}")


def extract_text_from_xlsx(file_content: bytes) -> str:
    """Extract text from Excel spreadsheet (.xlsx)."""
    try:
        wb = load_workbook(io.BytesIO(file_content), data_only=True)
        text_parts = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_text = [f"--- Sheet: {sheet_name} ---"]
            for row in sheet.iter_rows(values_only=True):
                row_values = [str(cell) if cell is not None else "" for cell in row]
                if any(v.strip() for v in row_values):
                    sheet_text.append(" | ".join(row_values))
            if len(sheet_text) > 1:
                text_parts.append("\n".join(sheet_text))
        return "\n\n".join(text_parts)
    except Exception as e:
        raise FileProcessingError(f"Failed to extract text from Excel: {str(e)}")


def extract_text_from_image(file_content: bytes, filename: str) -> str:
    """Extract text from image using Vision OCR (preferred) or Tesseract fallback."""
    # Try Vision OCR first (much better for math/formulas/diagrams)
    vision_parts = _ocr_images_with_vision([file_content])
    non_empty = [p for p in vision_parts if p]
    if non_empty:
        return "\n\n".join(non_empty)

    # Fallback to Tesseract
    if not OCR_AVAILABLE:
        raise FileProcessingError(
            "OCR is not available. Please install Tesseract and pytesseract, "
            "or configure ANTHROPIC_API_KEY for Vision OCR."
        )

    try:
        image = Image.open(io.BytesIO(file_content))
        # Convert to RGB if necessary (for PNG with alpha channel, etc.)
        if image.mode in ('RGBA', 'LA', 'P'):
            image = image.convert('RGB')
        text = pytesseract.image_to_string(image)
        if not text.strip():
            return f"[Image: {filename} - No text detected via OCR]"
        return text.strip()
    except Exception as e:
        raise FileProcessingError(f"Failed to extract text from image: {str(e)}")


def extract_text_from_text_file(file_content: bytes, filename: str) -> str:
    """Extract text from plain text files."""
    try:
        # Try common encodings
        for encoding in ['utf-8', 'utf-16', 'latin-1', 'cp1252']:
            try:
                return file_content.decode(encoding)
            except UnicodeDecodeError:
                continue
        raise FileProcessingError("Unable to decode text file with any supported encoding")
    except Exception as e:
        raise FileProcessingError(f"Failed to read text file: {str(e)}")


def extract_text_from_zip(file_content: bytes, filename: str) -> str:
    """Extract and process all supported files from a ZIP archive recursively."""
    logger.info(f"Processing ZIP archive: {filename}")
    try:
        extracted_texts = []
        total_size = 0
        files_processed = 0
        files_skipped = 0

        with zipfile.ZipFile(io.BytesIO(file_content), 'r') as zip_file:
            file_count = len([f for f in zip_file.infolist() if not f.is_dir()])
            logger.debug(f"ZIP contains {file_count} files")

            for zip_info in zip_file.infolist():
                # Skip directories
                if zip_info.is_dir():
                    continue

                # Check file extension
                file_ext = Path(zip_info.filename).suffix.lower()
                if file_ext not in SUPPORTED_EXTENSIONS or file_ext == '.zip':
                    # Skip unsupported files and nested zips (to prevent zip bombs)
                    logger.debug(f"Skipping unsupported file in ZIP: {zip_info.filename}")
                    files_skipped += 1
                    continue

                # Check cumulative size
                total_size += zip_info.file_size
                if total_size > MAX_FILE_SIZE:
                    logger.warning(f"ZIP extraction exceeded size limit at {total_size / (1024*1024):.2f} MB")
                    raise FileProcessingError(
                        f"Total extracted size exceeds {MAX_FILE_SIZE // (1024*1024)} MB limit"
                    )

                # Extract and process file
                inner_content = zip_file.read(zip_info.filename)
                try:
                    inner_text = process_file(inner_content, zip_info.filename)
                    if inner_text.strip():
                        extracted_texts.append(f"=== {zip_info.filename} ===\n{inner_text}")
                        files_processed += 1
                except FileProcessingError as e:
                    logger.debug(f"Failed to process file in ZIP: {zip_info.filename} - {str(e)}")
                    files_skipped += 1
                    continue

        logger.info(f"ZIP processing complete: {files_processed} files processed, {files_skipped} skipped")

        if not extracted_texts:
            raise FileProcessingError("No supported files found in ZIP archive")

        return "\n\n".join(extracted_texts)
    except zipfile.BadZipFile:
        logger.error(f"Invalid ZIP file: {filename}")
        raise FileProcessingError("Invalid or corrupted ZIP file")
    except Exception as e:
        if isinstance(e, FileProcessingError):
            raise
        logger.error(f"ZIP processing failed: {str(e)}")
        raise FileProcessingError(f"Failed to process ZIP archive: {str(e)}")


def _compress_image(image_bytes: bytes, max_width: int = 800) -> tuple[bytes, str]:
    """Resize image to max_width and compress. Returns (compressed_bytes, media_type)."""
    img = Image.open(io.BytesIO(image_bytes))
    if img.width > max_width:
        ratio = max_width / img.width
        img = img.resize((max_width, int(img.height * ratio)), Image.LANCZOS)
    # Keep PNG for images with transparency
    if img.mode in ('RGBA', 'LA', 'P'):
        buf = io.BytesIO()
        img.save(buf, format='PNG', optimize=True)
        return buf.getvalue(), 'image/png'
    else:
        # Convert to RGB for JPEG
        if img.mode != 'RGB':
            img = img.convert('RGB')
        buf = io.BytesIO()
        img.save(buf, format='JPEG', quality=85)
        return buf.getvalue(), 'image/jpeg'


def _extract_images_from_pptx(file_content: bytes) -> list[dict]:
    """Extract images from a PowerPoint file with slide context."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    images: list[dict] = []
    try:
        prs = Presentation(io.BytesIO(file_content))
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_title = ""
            if slide.shapes.title:
                slide_title = slide.shapes.title.text or ""
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        img_blob = shape.image.blob
                        img_content_type = shape.image.content_type
                        context = f"Slide {slide_num}"
                        if slide_title:
                            context += f": {slide_title}"
                        images.append({
                            "image_bytes": img_blob,
                            "content_type": img_content_type,
                            "position_context": context,
                        })
                    except Exception as e:
                        logger.debug(f"Failed to extract image from slide {slide_num}: {e}")
    except Exception as e:
        logger.debug(f"Failed to extract images from pptx: {e}")
    return images


def _extract_images_from_pdf(file_content: bytes) -> list[dict]:
    """Extract images from a PDF file with page context.

    First tries PyPDF2 to extract embedded raster images. If none found,
    falls back to PyMuPDF (fitz) to render each page as a PNG, capturing
    vector graphics (diagrams, charts, math figures) that PyPDF2 cannot extract.
    """
    images: list[dict] = []
    try:
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
        for page_num, page in enumerate(pdf_reader.pages, 1):
            try:
                if hasattr(page, 'images') and page.images:
                    for img in page.images:
                        images.append({
                            "image_bytes": img.data,
                            "content_type": None,  # Will be detected later
                            "position_context": f"Page {page_num}",
                        })
                elif '/Resources' in page and '/XObject' in page['/Resources']:
                    x_objects = page['/Resources']['/XObject'].get_object()
                    for obj_name in x_objects:
                        obj = x_objects[obj_name].get_object()
                        if obj.get('/Subtype') == '/Image':
                            try:
                                data = obj.get_data()
                                images.append({
                                    "image_bytes": data,
                                    "content_type": None,
                                    "position_context": f"Page {page_num}",
                                })
                            except Exception as e:
                                logger.debug(f"Failed to extract image {obj_name} from page {page_num}: {e}")
            except Exception as e:
                logger.debug(f"Failed to extract images from PDF page {page_num}: {e}")
    except Exception as e:
        logger.debug(f"Failed to extract images from PDF: {e}")

    # Fallback: if PyPDF2 found no raster images, render pages as images
    # using PyMuPDF to capture vector graphics (diagrams, math figures, etc.)
    if not images:
        images = _render_pdf_pages_as_images(file_content)

    return images


def _render_pdf_pages_as_images(file_content: bytes) -> list[dict]:
    """Render PDF pages as PNG images using PyMuPDF (fitz).

    This captures vector graphics (drawn shapes, diagrams, math figures)
    that cannot be extracted as embedded raster images by PyPDF2.
    Renders at 150 DPI for a good balance of quality vs file size.
    """
    images: list[dict] = []
    try:
        import fitz  # PyMuPDF
    except ImportError:
        logger.warning("PyMuPDF (fitz) not installed — cannot render PDF pages as images")
        return images

    try:
        doc = fitz.open(stream=file_content, filetype="pdf")
        # Cap at 20 pages to avoid excessive processing
        max_pages = min(len(doc), _MAX_IMAGES_PER_DOC)
        for page_num in range(max_pages):
            try:
                page = doc[page_num]
                # Render at 150 DPI (default is 72 DPI, so zoom = 150/72 ≈ 2.08)
                zoom = 150 / 72
                mat = fitz.Matrix(zoom, zoom)
                pix = page.get_pixmap(matrix=mat)
                png_bytes = pix.tobytes("png")
                images.append({
                    "image_bytes": png_bytes,
                    "content_type": "image/png",
                    "position_context": f"Page {page_num + 1}",
                })
            except Exception as e:
                logger.debug(f"Failed to render PDF page {page_num + 1} as image: {e}")
        doc.close()
        logger.info(f"Rendered {len(images)} PDF pages as images via PyMuPDF fallback")
    except Exception as e:
        logger.debug(f"PyMuPDF PDF rendering failed: {e}")
    return images


def _extract_docx_images_with_context(file_content: bytes) -> list[dict]:
    """Extract images from a DOCX with surrounding paragraph context.

    Walks the document XML to find inline/anchor images and associates them
    with the paragraphs before and after each image.
    """
    images: list[dict] = []
    try:
        doc = WordDocument(io.BytesIO(file_content))
        # Build a map of relationship IDs to image names for the main document part
        rels = doc.part.rels
        rid_to_image: dict[str, str] = {}
        for rel_id, rel in rels.items():
            if "image" in (rel.reltype or ""):
                rid_to_image[rel_id] = rel.target_ref

        # Extract image blobs from the ZIP
        image_blobs: dict[str, bytes] = {}
        with zipfile.ZipFile(io.BytesIO(file_content), "r") as zf:
            for name in zf.namelist():
                if name.startswith("word/media/") and Path(name).suffix.lower() in IMAGE_EXTENSIONS:
                    image_blobs[name] = zf.read(name)

        # Walk paragraphs to find images and capture surrounding text
        paragraphs = doc.paragraphs
        para_texts = [p.text.strip() for p in paragraphs]

        # Namespace for finding image references
        ns_drawing = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
        ns_blip = "http://schemas.openxmlformats.org/drawingml/2006/main"
        ns_r = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

        for para_idx, para in enumerate(paragraphs):
            # Find all blip elements (inline and anchor images) in this paragraph's XML
            blips = para._element.findall(f".//{{{ns_blip}}}blip")
            for blip in blips:
                embed_rid = blip.get(f"{{{ns_r}}}embed")
                if not embed_rid or embed_rid not in rid_to_image:
                    continue
                target = rid_to_image[embed_rid]
                # Target may be relative like "media/image1.png" — normalize
                if not target.startswith("word/"):
                    target = "word/" + target
                if target not in image_blobs:
                    continue

                # Build context from surrounding paragraphs
                context_parts = []
                if para_idx > 0 and para_texts[para_idx - 1]:
                    context_parts.append(para_texts[para_idx - 1])
                if para_texts[para_idx]:
                    context_parts.append(para_texts[para_idx])
                if para_idx + 1 < len(para_texts) and para_texts[para_idx + 1]:
                    context_parts.append(para_texts[para_idx + 1])
                context = " | ".join(context_parts) if context_parts else None

                images.append({
                    "image_bytes": image_blobs[target],
                    "content_type": None,
                    "position_context": context,
                })

        # If XML walking found nothing, fall back to extracting all images
        if not images:
            raw_images = _extract_images_from_docx(file_content)
            for img_bytes in raw_images:
                images.append({
                    "image_bytes": img_bytes,
                    "content_type": None,
                    "position_context": None,
                })

    except Exception as e:
        logger.debug(f"DOCX image extraction with context failed, falling back: {e}")
        # Fallback: use the simple extraction
        raw_images = _extract_images_from_docx(file_content)
        for img_bytes in raw_images:
            images.append({
                "image_bytes": img_bytes,
                "content_type": None,
                "position_context": None,
            })
    return images


# Maximum number of images to extract per document
_MAX_IMAGES_PER_DOC = 20


def extract_images_from_file(file_content: bytes, filename: str) -> list[dict]:
    """Extract images from a document file.

    Returns list of dicts with keys:
        image_data: bytes (compressed, max 800px width)
        media_type: str
        description: str | None (from Vision OCR)
        position_context: str | None (nearby text)
        position_index: int
        file_size: int
    """
    ext = Path(filename).suffix.lower()
    raw_images: list[dict] = []

    if ext in ('.docx',):
        raw_images = _extract_docx_images_with_context(file_content)
    elif ext == '.pdf':
        raw_images = _extract_images_from_pdf(file_content)
    elif ext == '.pptx':
        raw_images = _extract_images_from_pptx(file_content)
    elif ext in IMAGE_EXTENSIONS:
        # Single image file — the file itself is the image
        raw_images = [{
            "image_bytes": file_content,
            "content_type": None,
            "position_context": None,
        }]
    else:
        # Unsupported format for image extraction
        return []

    if not raw_images:
        return []

    # Filter out tiny images (< 1KB)
    raw_images = [img for img in raw_images if len(img["image_bytes"]) > 1024]
    if not raw_images:
        return []

    # Cap at _MAX_IMAGES_PER_DOC — keep the largest if over limit
    if len(raw_images) > _MAX_IMAGES_PER_DOC:
        raw_images.sort(key=lambda x: len(x["image_bytes"]), reverse=True)
        raw_images = raw_images[:_MAX_IMAGES_PER_DOC]

    # Compress images and collect bytes for OCR
    compressed: list[tuple[bytes, str]] = []
    for img_info in raw_images:
        try:
            comp_bytes, comp_media = _compress_image(img_info["image_bytes"])
            compressed.append((comp_bytes, comp_media))
        except Exception as e:
            logger.debug(f"Image compression failed, using original: {e}")
            media = img_info.get("content_type") or _get_image_media_type(img_info["image_bytes"])
            compressed.append((img_info["image_bytes"], media))

    # Run Vision OCR for descriptions (use compressed images)
    ocr_images = [comp[0] for comp in compressed]
    descriptions = _ocr_images_with_vision(ocr_images)

    # Build result dicts
    results: list[dict] = []
    for idx, (img_info, (comp_bytes, comp_media)) in enumerate(zip(raw_images, compressed)):
        desc = descriptions[idx] if idx < len(descriptions) else ""
        results.append({
            "image_data": comp_bytes,
            "media_type": comp_media,
            "description": desc if desc else None,
            "position_context": img_info.get("position_context"),
            "position_index": idx,
            "file_size": len(comp_bytes),
        })

    return results


def process_file(file_content: bytes, filename: str) -> str:
    """
    Process a file and extract its text content.

    Args:
        file_content: Raw bytes of the file
        filename: Original filename (used to determine file type)

    Returns:
        Extracted text content

    Raises:
        FileProcessingError: If file cannot be processed
    """
    logger.info(f"Processing file: {filename}")
    validate_file(file_content, filename)

    ext = Path(filename).suffix.lower()

    if ext == '.pdf':
        return extract_text_from_pdf(file_content)

    elif ext in ('.doc', '.docx'):
        if ext == '.doc':
            raise FileProcessingError(
                "Legacy .doc format is not supported. Please convert to .docx"
            )
        return extract_text_from_docx(file_content)

    elif ext == '.pptx':
        return extract_text_from_pptx(file_content)

    elif ext == '.ppt':
        raise FileProcessingError(
            "Legacy .ppt format is not supported. Please convert to .pptx"
        )

    elif ext in ('.xlsx', '.xls'):
        if ext == '.xls':
            raise FileProcessingError(
                "Legacy .xls format is not supported. Please convert to .xlsx"
            )
        return extract_text_from_xlsx(file_content)

    elif ext == '.csv':
        return extract_text_from_text_file(file_content, filename)

    elif ext in IMAGE_EXTENSIONS:
        return extract_text_from_image(file_content, filename)

    elif ext in ('.txt', '.md', '.rtf'):
        return extract_text_from_text_file(file_content, filename)

    elif ext == '.zip':
        return extract_text_from_zip(file_content, filename)

    else:
        raise FileProcessingError(f"Unsupported file type: {ext}")


async def process_uploaded_file(file: BinaryIO, filename: str) -> str:
    """
    Async wrapper for processing uploaded files.

    Args:
        file: File-like object with read() method
        filename: Original filename

    Returns:
        Extracted text content
    """
    file_content = file.read()
    return process_file(file_content, filename)


def get_supported_formats() -> dict:
    """Return information about supported file formats."""
    return {
        "documents": [".pdf", ".docx", ".txt", ".md", ".rtf"],
        "spreadsheets": [".xlsx", ".csv"],
        "presentations": [".pptx"],
        "images": list(IMAGE_EXTENSIONS),
        "archives": [".zip"],
        "max_file_size_mb": MAX_FILE_SIZE // (1024 * 1024),
        "ocr_available": OCR_AVAILABLE,
        "scanned_pdf_ocr": OCR_AVAILABLE and PDF2IMAGE_AVAILABLE,
    }
